import copy
import sys
import six
sys.path.append('/Users/dave/projects/thm-dev/python-pptx/')
from pptx import Presentation

nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

prs = Presentation('ppt-with-transitions.pptx')
output_prs = Presentation()


def _get_blank_slide_layout(pres):
    layout_items_count = [
        len(layout.placeholders) for layout in pres.slide_layouts
    ]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def copy_slide(pres, pres1, index, ignored_shape_ids=[]):
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres1.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        if shp.id in ignored_shape_ids:
            continue
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for key, value in six.iteritems(source.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    return dest


def process_slide(slide_index):
    timing = prs.slides[slide_index].element.timing
    target_ids_with_click_effects = {
        int(e.attrib['spid']) for e in
        timing.xpath(
            '//p:cTn[@nodeType="clickEffect"]//p:spTgt',
            namespaces=nsmap,
        )
    }
    shapes = prs.slides[slide_index].shapes
    shapes_with_animations = [
        shape for shape in shapes
        if shape.id in target_ids_with_click_effects
    ]
    shape_ids_with_animations = {shape.id for shape in shapes_with_animations}

    for i, slide in enumerate(prs.slides):
        # Remake the presentation since you can't add a slide to an existing
        if i == slide_index:
            copy_slide(prs, output_prs, i, shape_ids_with_animations)
        copy_slide(prs, output_prs, i)
    output_prs.save('prs-extended.pptx')

if __name__ == '__main__':
    process_slide(0)

# import ipdb;ipdb.set_trace()
